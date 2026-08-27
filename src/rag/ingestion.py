"""Document ingestion for PDF and PPTX files."""
from pathlib import Path
from typing import Generator
from pypdf import PdfReader
from pptx import Presentation
import re


def extract_text_from_pdf(file_path: Path) -> str:
    """Extract text from PDF file."""
    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text)
    return "\n\n".join(texts)


def extract_text_from_pptx(file_path: Path) -> str:
    """Extract text from PPTX file."""
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
        if slide_texts:
            texts.append("\n".join(slide_texts))
    return "\n\n---\n\n".join(texts)


def extract_text(file_path: Path) -> str:
    """Extract text based on file extension."""
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    elif suffix == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        raise ValueError(f"Unsupported file format: {suffix}")


def iter_documents(docs_dir: Path) -> Generator[tuple[str, str], None, None]:
    """Yield (filename, text) for each supported document in directory."""
    for file_path in docs_dir.iterdir():
        if file_path.suffix.lower() in {".pdf", ".pptx"}:
            try:
                text = extract_text(file_path)
                if text.strip():
                    yield file_path.name, text
            except Exception as e:
                print(f"Error processing {file_path.name}: {e}")


def clean_text(text: str) -> str:
    """Clean extracted text."""
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"[ \t]+", " ", text)
    return text.strip()